import io

from docx import Document
from pypdf import PdfReader
from pptx import Presentation

LEGACY_DOC_MESSAGE = "Legacy .doc files are not supported. Please save as .docx and try again."


def extract_text_from_pdf(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    pages = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text.strip())
    return "\n\n".join(pages)


def extract_text_from_pptx(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
        if lines:
            slides.append(f"[Slide {i}]\n" + "\n".join(lines))
    return "\n\n".join(slides)


def extract_text_from_docx(data: bytes) -> str:
    doc = Document(io.BytesIO(data))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def extract_text(data: bytes, filename: str) -> str:
    lower = filename.lower()
    if lower.endswith(".pdf"):
        return extract_text_from_pdf(data)
    if lower.endswith(".pptx"):
        return extract_text_from_pptx(data)
    if lower.endswith(".docx"):
        return extract_text_from_docx(data)
    if lower.endswith(".doc"):
        raise ValueError(LEGACY_DOC_MESSAGE)
    if lower.endswith(".txt") or lower.endswith(".md"):
        return data.decode("utf-8", errors="replace")
    raise ValueError(
        "Unsupported file type. Please upload a PDF, PPTX, TXT, MD, or DOCX file."
    )